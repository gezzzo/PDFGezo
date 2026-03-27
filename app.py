import os
import uuid
import shutil
import zipfile
import subprocess
from flask import Flask, render_template, request, send_file, jsonify
from pdf2image import convert_from_path
from PyPDF2 import PdfMerger, PdfReader, PdfWriter

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50 MB limit

UPLOAD_DIR = os.path.join(os.path.dirname(__file__), 'uploads')
OUTPUT_DIR = os.path.join(os.path.dirname(__file__), 'outputs')

os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)


# ─── Pages ────────────────────────────────────────────────────────────
@app.route('/')
def index():
    return render_template('index.html')


@app.route('/pdf-to-images')
def pdf_to_images_page():
    return render_template('pdf_to_images.html')


@app.route('/merge-pdf')
def merge_pdf_page():
    return render_template('merge_pdf.html')


@app.route('/split-pdf')
def split_pdf_page():
    return render_template('split_pdf.html')


@app.route('/compress-pdf')
def compress_pdf_page():
    return render_template('compress_pdf.html')


@app.route('/jpg-to-pdf')
def jpg_to_pdf_page():
    return render_template('jpg_to_pdf.html')


@app.route('/word-to-pdf')
def word_to_pdf_page():
    return render_template('word_to_pdf.html')


@app.route('/pptx-to-pdf')
def pptx_to_pdf_page():
    return render_template('pptx_to_pdf.html')

@app.route('/excel-to-pdf')
def excel_to_pdf_page():
    return render_template('excel_to_pdf.html')

@app.route('/html-to-pdf')
def html_to_pdf_page():
    return render_template('html_to_pdf.html')

@app.route('/csv-to-pdf')
def csv_to_pdf_page():
    return render_template('csv_to_pdf.html')

@app.route('/pdf-to-word')
def pdf_to_word_page():
    return render_template('pdf_to_word.html')

@app.route('/pdf-to-pptx')
def pdf_to_pptx_page():
    return render_template('pdf_to_pptx.html')

@app.route('/pdf-to-excel')
def pdf_to_excel_page():
    return render_template('pdf_to_excel.html')

@app.route('/pdf-to-csv')
def pdf_to_csv_page():
    return render_template('pdf_to_csv.html')


# ─── API ──────────────────────────────────────────────────────────────
@app.route('/api/pdf-to-images', methods=['POST'])
def api_pdf_to_images():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400

    if not file.filename.lower().endswith('.pdf'):
        return jsonify({'error': 'Only PDF files are allowed'}), 400

    # Create unique job id
    job_id = str(uuid.uuid4())
    job_upload = os.path.join(UPLOAD_DIR, job_id)
    job_output = os.path.join(OUTPUT_DIR, job_id)
    os.makedirs(job_upload, exist_ok=True)
    os.makedirs(job_output, exist_ok=True)

    pdf_path = os.path.join(job_upload, file.filename)
    file.save(pdf_path)

    try:
        # Convert PDF pages → images
        dpi = int(request.form.get('dpi', 200))
        fmt = request.form.get('format', 'png').lower()
        if fmt not in ('png', 'jpg', 'jpeg'):
            fmt = 'png'

        images = convert_from_path(pdf_path, dpi=dpi)

        base_name = os.path.splitext(file.filename)[0]

        image_paths = []
        for i, img in enumerate(images, start=1):
            img_name = f"{base_name}_page_{i}.{fmt}"
            img_path = os.path.join(job_output, img_name)
            img.save(img_path, fmt.upper().replace('JPG', 'JPEG'))
            image_paths.append(img_path)

        # Bundle into ZIP
        zip_name = f"{base_name}_images.zip"
        zip_path = os.path.join(job_output, zip_name)
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for p in image_paths:
                zf.write(p, os.path.basename(p))

        response = send_file(
            zip_path,
            mimetype='application/zip',
            as_attachment=True,
            download_name=zip_name
        )

        # Cleanup after response is sent
        @response.call_on_close
        def _cleanup():
            shutil.rmtree(job_upload, ignore_errors=True)
            shutil.rmtree(job_output, ignore_errors=True)

        return response

    except Exception as e:
        # Cleanup on error
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': str(e)}), 500


# ─── Merge PDF API ───────────────────────────────────────────────────
@app.route('/api/merge-pdf', methods=['POST'])
def api_merge_pdf():
    files = request.files.getlist('files')
    if not files or len(files) < 2:
        return jsonify({'error': 'Please upload at least 2 PDF files.'}), 400

    for f in files:
        if not f.filename.lower().endswith('.pdf'):
            return jsonify({'error': f'File "{f.filename}" is not a PDF.'}), 400

    job_id = str(uuid.uuid4())
    job_upload = os.path.join(UPLOAD_DIR, job_id)
    job_output = os.path.join(OUTPUT_DIR, job_id)
    os.makedirs(job_upload, exist_ok=True)
    os.makedirs(job_output, exist_ok=True)

    try:
        saved_paths = []
        for f in files:
            path = os.path.join(job_upload, f.filename)
            f.save(path)
            saved_paths.append(path)

        merger = PdfMerger()
        for p in saved_paths:
            merger.append(p)

        merged_name = 'merged.pdf'
        merged_path = os.path.join(job_output, merged_name)
        merger.write(merged_path)
        merger.close()

        response = send_file(
            merged_path,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=merged_name
        )

        @response.call_on_close
        def _cleanup():
            shutil.rmtree(job_upload, ignore_errors=True)
            shutil.rmtree(job_output, ignore_errors=True)

        return response

    except Exception as e:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': str(e)}), 500


# ─── Split PDF API ───────────────────────────────────────────────────
@app.route('/api/split-pdf', methods=['POST'])
def api_split_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400

    file = request.files['file']
    if not file.filename.lower().endswith('.pdf'):
        return jsonify({'error': 'Only PDF files are allowed'}), 400

    ranges_str = request.form.get('ranges', '').strip()
    if not ranges_str:
        return jsonify({'error': 'Please specify page ranges.'}), 400

    job_id = str(uuid.uuid4())
    job_upload = os.path.join(UPLOAD_DIR, job_id)
    job_output = os.path.join(OUTPUT_DIR, job_id)
    os.makedirs(job_upload, exist_ok=True)
    os.makedirs(job_output, exist_ok=True)

    pdf_path = os.path.join(job_upload, file.filename)
    file.save(pdf_path)

    try:
        reader = PdfReader(pdf_path)
        total_pages = len(reader.pages)
        base_name = os.path.splitext(file.filename)[0]

        # Parse ranges like "1-3,5,7-9"
        pages_to_extract = []
        for part in ranges_str.split(','):
            part = part.strip()
            if '-' in part:
                start, end = part.split('-', 1)
                s, e = int(start.strip()), int(end.strip())
                if s < 1 or e > total_pages or s > e:
                    return jsonify({'error': f'Invalid range {part}. PDF has {total_pages} pages.'}), 400
                pages_to_extract.extend(range(s, e + 1))
            else:
                p = int(part)
                if p < 1 or p > total_pages:
                    return jsonify({'error': f'Page {p} is out of range. PDF has {total_pages} pages.'}), 400
                pages_to_extract.append(p)

        if not pages_to_extract:
            return jsonify({'error': 'No valid pages specified.'}), 400

        # Remove duplicates and sort
        pages_to_extract = sorted(set(pages_to_extract))

        writer = PdfWriter()
        for p in pages_to_extract:
            writer.add_page(reader.pages[p - 1])

        out_name = f"{base_name}_split.pdf"
        out_path = os.path.join(job_output, out_name)
        with open(out_path, 'wb') as f:
            writer.write(f)

        response = send_file(
            out_path,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=out_name
        )

        @response.call_on_close
        def _cleanup():
            shutil.rmtree(job_upload, ignore_errors=True)
            shutil.rmtree(job_output, ignore_errors=True)

        return response

    except ValueError:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': 'Invalid page range format. Use format like: 1-3,5,7-9'}), 400
    except Exception as e:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': str(e)}), 500


# ─── Compress PDF API ────────────────────────────────────────────────
@app.route('/api/compress-pdf', methods=['POST'])
def api_compress_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400

    file = request.files['file']
    if not file.filename.lower().endswith('.pdf'):
        return jsonify({'error': 'Only PDF files are allowed'}), 400

    quality = request.form.get('quality', 'medium').lower()
    gs_settings = {
        'low':    '/screen',      # 72 DPI — smallest size
        'medium': '/ebook',       # 150 DPI — good balance
        'high':   '/printer',     # 300 DPI — high quality
    }
    pdf_setting = gs_settings.get(quality, '/ebook')

    job_id = str(uuid.uuid4())
    job_upload = os.path.join(UPLOAD_DIR, job_id)
    job_output = os.path.join(OUTPUT_DIR, job_id)
    os.makedirs(job_upload, exist_ok=True)
    os.makedirs(job_output, exist_ok=True)

    pdf_path = os.path.join(job_upload, file.filename)
    file.save(pdf_path)

    try:
        base_name = os.path.splitext(file.filename)[0]
        out_name = f"{base_name}_compressed.pdf"
        out_path = os.path.join(job_output, out_name)

        original_size = os.path.getsize(pdf_path)

        # Use Ghostscript for proper PDF compression
        gs_cmd = [
            'gs', '-sDEVICE=pdfwrite', '-dCompatibilityLevel=1.4',
            f'-dPDFSETTINGS={pdf_setting}',
            '-dNOPAUSE', '-dQUIET', '-dBATCH',
            f'-sOutputFile={out_path}', pdf_path
        ]
        result = subprocess.run(gs_cmd, capture_output=True, text=True, timeout=120)

        if result.returncode != 0:
            return jsonify({'error': 'Compression failed. ' + result.stderr[:200]}), 500

        compressed_size = os.path.getsize(out_path)
        ratio = round((1 - compressed_size / original_size) * 100, 1) if original_size > 0 else 0

        response = send_file(
            out_path,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=out_name
        )
        response.headers['X-Original-Size'] = str(original_size)
        response.headers['X-Compressed-Size'] = str(compressed_size)
        response.headers['X-Compression-Ratio'] = str(ratio)

        @response.call_on_close
        def _cleanup():
            shutil.rmtree(job_upload, ignore_errors=True)
            shutil.rmtree(job_output, ignore_errors=True)

        return response

    except subprocess.TimeoutExpired:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': 'Compression timed out. Try a smaller file.'}), 500
    except Exception as e:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': str(e)}), 500


# ─── Images to PDF API ──────────────────────────────────────────────────
@app.route('/api/jpg-to-pdf', methods=['POST'])
def api_jpg_to_pdf():
    files = request.files.getlist('files')
    if not files or len(files) < 1:
        return jsonify({'error': 'Please upload at least 1 image.'}), 400

    allowed_ext = {'.jpg', '.jpeg', '.png', '.webp'}
    for f in files:
        ext = os.path.splitext(f.filename)[0]
        if os.path.splitext(f.filename)[1].lower() not in allowed_ext:
            return jsonify({'error': f'File "{f.filename}" is not a supported image (JPG, PNG, WEBP).'}), 400

    job_id = str(uuid.uuid4())
    job_upload = os.path.join(UPLOAD_DIR, job_id)
    job_output = os.path.join(OUTPUT_DIR, job_id)
    os.makedirs(job_upload, exist_ok=True)
    os.makedirs(job_output, exist_ok=True)

    try:
        from PIL import Image

        images = []
        for f in files:
            path = os.path.join(job_upload, f.filename)
            f.save(path)
            img = Image.open(path)
            if img.mode in ('RGBA', 'P', 'LA'):
                img = img.convert('RGB')
            elif img.mode != 'RGB':
                img = img.convert('RGB')
            images.append(img)

        out_name = 'images_converted.pdf'
        out_path = os.path.join(job_output, out_name)

        if len(images) == 1:
            images[0].save(out_path, 'PDF')
        else:
            images[0].save(out_path, 'PDF', save_all=True, append_images=images[1:])

        response = send_file(
            out_path,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=out_name
        )

        @response.call_on_close
        def _cleanup():
            shutil.rmtree(job_upload, ignore_errors=True)
            shutil.rmtree(job_output, ignore_errors=True)

        return response

    except Exception as e:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': str(e)}), 500


# ─── Word to PDF API ─────────────────────────────────────────────────
@app.route('/api/word-to-pdf', methods=['POST'])
def api_word_to_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded.'}), 400

    file = request.files['file']
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in ('.doc', '.docx'):
        return jsonify({'error': 'Only .doc and .docx files are supported.'}), 400

    job_id = str(uuid.uuid4())
    job_upload = os.path.join(UPLOAD_DIR, job_id)
    job_output = os.path.join(OUTPUT_DIR, job_id)
    os.makedirs(job_upload, exist_ok=True)
    os.makedirs(job_output, exist_ok=True)

    try:
        input_path = os.path.join(job_upload, file.filename)
        file.save(input_path)

        # Use LibreOffice headless to convert
        result = subprocess.run([
            'soffice', '--headless', '--convert-to', 'pdf',
            '--outdir', job_output, input_path
        ], capture_output=True, timeout=120)

        if result.returncode != 0:
            return jsonify({'error': 'Conversion failed. Please try again.'}), 500

        # Find the output PDF
        base_name = os.path.splitext(file.filename)[0] + '.pdf'
        out_path = os.path.join(job_output, base_name)

        if not os.path.exists(out_path):
            # Try to find any PDF in output dir
            pdfs = [f for f in os.listdir(job_output) if f.endswith('.pdf')]
            if pdfs:
                out_path = os.path.join(job_output, pdfs[0])
                base_name = pdfs[0]
            else:
                return jsonify({'error': 'Conversion produced no output.'}), 500

        response = send_file(
            out_path,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=base_name
        )

        @response.call_on_close
        def _cleanup():
            shutil.rmtree(job_upload, ignore_errors=True)
            shutil.rmtree(job_output, ignore_errors=True)

        return response

    except subprocess.TimeoutExpired:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': 'Conversion timed out. The file may be too large.'}), 500
    except Exception as e:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': str(e)}), 500


# ─── PowerPoint to PDF API ─────────────────────────────────────────────
@app.route('/api/pptx-to-pdf', methods=['POST'])
def api_pptx_to_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded.'}), 400

    file = request.files['file']
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in ('.ppt', '.pptx'):
        return jsonify({'error': 'Only .ppt and .pptx files are supported.'}), 400

    job_id = str(uuid.uuid4())
    job_upload = os.path.join(UPLOAD_DIR, job_id)
    job_output = os.path.join(OUTPUT_DIR, job_id)
    os.makedirs(job_upload, exist_ok=True)
    os.makedirs(job_output, exist_ok=True)

    try:
        input_path = os.path.join(job_upload, file.filename)
        file.save(input_path)

        # Use LibreOffice headless to convert
        result = subprocess.run([
            'soffice', '--headless', '--convert-to', 'pdf',
            '--outdir', job_output, input_path
        ], capture_output=True, timeout=120)

        if result.returncode != 0:
            return jsonify({'error': 'Conversion failed. Please try again.'}), 500

        # Find the output PDF
        base_name = os.path.splitext(file.filename)[0] + '.pdf'
        out_path = os.path.join(job_output, base_name)

        if not os.path.exists(out_path):
            # Try to find any PDF in output dir
            pdfs = [f for f in os.listdir(job_output) if f.endswith('.pdf')]
            if pdfs:
                out_path = os.path.join(job_output, pdfs[0])
                base_name = pdfs[0]
            else:
                return jsonify({'error': 'Conversion produced no output.'}), 500

        response = send_file(
            out_path,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=base_name
        )

        @response.call_on_close
        def _cleanup():
            shutil.rmtree(job_upload, ignore_errors=True)
            shutil.rmtree(job_output, ignore_errors=True)

        return response

    except subprocess.TimeoutExpired:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': 'Conversion timed out. The file may be too large.'}), 500
    except Exception as e:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': str(e)}), 500


# ─── Excel to PDF API ────────────────────────────────────────────────
@app.route('/api/excel-to-pdf', methods=['POST'])
def api_excel_to_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded.'}), 400

    file = request.files['file']
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in ('.xls', '.xlsx'):
        return jsonify({'error': 'Only .xls and .xlsx files are supported.'}), 400

    job_id = str(uuid.uuid4())
    job_upload = os.path.join(UPLOAD_DIR, job_id)
    job_output = os.path.join(OUTPUT_DIR, job_id)
    os.makedirs(job_upload, exist_ok=True)
    os.makedirs(job_output, exist_ok=True)

    try:
        input_path = os.path.join(job_upload, file.filename)
        file.save(input_path)

        # Use LibreOffice headless to convert
        result = subprocess.run([
            'soffice', '--headless', '--convert-to', 'pdf',
            '--outdir', job_output, input_path
        ], capture_output=True, timeout=120)

        if result.returncode != 0:
            return jsonify({'error': 'Conversion failed. Please try again.'}), 500

        # Find the output PDF
        base_name = os.path.splitext(file.filename)[0] + '.pdf'
        out_path = os.path.join(job_output, base_name)

        if not os.path.exists(out_path):
            pdfs = [f for f in os.listdir(job_output) if f.endswith('.pdf')]
            if pdfs:
                out_path = os.path.join(job_output, pdfs[0])
                base_name = pdfs[0]
            else:
                return jsonify({'error': 'Conversion produced no output.'}), 500

        response = send_file(
            out_path,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=base_name
        )

        @response.call_on_close
        def _cleanup():
            shutil.rmtree(job_upload, ignore_errors=True)
            shutil.rmtree(job_output, ignore_errors=True)

        return response

    except subprocess.TimeoutExpired:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': 'Conversion timed out. The file may be too large.'}), 500
    except Exception as e:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': str(e)}), 500


# ─── HTML to PDF API ─────────────────────────────────────────────────
@app.route('/api/html-to-pdf', methods=['POST'])
def api_html_to_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded.'}), 400

    file = request.files['file']
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in ('.html', '.htm'):
        return jsonify({'error': 'Only .html and .htm files are supported.'}), 400

    job_id = str(uuid.uuid4())
    job_upload = os.path.join(UPLOAD_DIR, job_id)
    job_output = os.path.join(OUTPUT_DIR, job_id)
    os.makedirs(job_upload, exist_ok=True)
    os.makedirs(job_output, exist_ok=True)

    try:
        input_path = os.path.join(job_upload, file.filename)
        file.save(input_path)

        # Use LibreOffice headless to convert
        result = subprocess.run([
            'soffice', '--headless', '--convert-to', 'pdf',
            '--outdir', job_output, input_path
        ], capture_output=True, timeout=120)

        if result.returncode != 0:
            return jsonify({'error': 'Conversion failed. Please try again.'}), 500

        # Find the output PDF
        base_name = os.path.splitext(file.filename)[0] + '.pdf'
        out_path = os.path.join(job_output, base_name)

        if not os.path.exists(out_path):
            pdfs = [f for f in os.listdir(job_output) if f.endswith('.pdf')]
            if pdfs:
                out_path = os.path.join(job_output, pdfs[0])
                base_name = pdfs[0]
            else:
                return jsonify({'error': 'Conversion produced no output.'}), 500

        response = send_file(
            out_path,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=base_name
        )

        @response.call_on_close
        def _cleanup():
            shutil.rmtree(job_upload, ignore_errors=True)
            shutil.rmtree(job_output, ignore_errors=True)

        return response

    except subprocess.TimeoutExpired:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': 'Conversion timed out. The file may be too large.'}), 500
    except Exception as e:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': str(e)}), 500


# ─── CSV to PDF API ──────────────────────────────────────────────────
@app.route('/api/csv-to-pdf', methods=['POST'])
def api_csv_to_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded.'}), 400

    file = request.files['file']
    ext = os.path.splitext(file.filename)[1].lower()
    if ext != '.csv':
        return jsonify({'error': 'Only .csv files are supported.'}), 400

    job_id = str(uuid.uuid4())
    job_upload = os.path.join(UPLOAD_DIR, job_id)
    job_output = os.path.join(OUTPUT_DIR, job_id)
    os.makedirs(job_upload, exist_ok=True)
    os.makedirs(job_output, exist_ok=True)

    try:
        input_path = os.path.join(job_upload, file.filename)
        file.save(input_path)

        # Use LibreOffice headless to convert CSV → PDF
        result = subprocess.run([
            'soffice', '--headless', '--convert-to', 'pdf',
            '--outdir', job_output, input_path
        ], capture_output=True, timeout=120)

        if result.returncode != 0:
            return jsonify({'error': 'Conversion failed. Please try again.'}), 500

        # Find the output PDF
        base_name = os.path.splitext(file.filename)[0] + '.pdf'
        out_path = os.path.join(job_output, base_name)

        if not os.path.exists(out_path):
            pdfs = [f for f in os.listdir(job_output) if f.endswith('.pdf')]
            if pdfs:
                out_path = os.path.join(job_output, pdfs[0])
                base_name = pdfs[0]
            else:
                return jsonify({'error': 'Conversion produced no output.'}), 500

        response = send_file(
            out_path,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=base_name
        )

        @response.call_on_close
        def _cleanup():
            shutil.rmtree(job_upload, ignore_errors=True)
            shutil.rmtree(job_output, ignore_errors=True)

        return response

    except subprocess.TimeoutExpired:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': 'Conversion timed out. The file may be too large.'}), 500
    except Exception as e:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': str(e)}), 500


# ─── PDF to Word API ─────────────────────────────────────────────────
@app.route('/api/pdf-to-word', methods=['POST'])
def api_pdf_to_word():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded.'}), 400

    file = request.files['file']
    ext = os.path.splitext(file.filename)[1].lower()
    if ext != '.pdf':
        return jsonify({'error': 'Only PDF files are supported.'}), 400

    job_id = str(uuid.uuid4())
    job_upload = os.path.join(UPLOAD_DIR, job_id)
    job_output = os.path.join(OUTPUT_DIR, job_id)
    os.makedirs(job_upload, exist_ok=True)
    os.makedirs(job_output, exist_ok=True)

    try:
        from pdf2docx import Converter

        input_path = os.path.join(job_upload, file.filename)
        file.save(input_path)

        base_name = os.path.splitext(file.filename)[0] + '.docx'
        out_path = os.path.join(job_output, base_name)

        cv = Converter(input_path)
        cv.convert(out_path)
        cv.close()

        if not os.path.exists(out_path):
            return jsonify({'error': 'Conversion produced no output.'}), 500

        response = send_file(
            out_path,
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            as_attachment=True,
            download_name=base_name
        )

        @response.call_on_close
        def _cleanup():
            shutil.rmtree(job_upload, ignore_errors=True)
            shutil.rmtree(job_output, ignore_errors=True)

        return response

    except Exception as e:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': str(e)}), 500


# ─── PDF to PowerPoint API ───────────────────────────────────────────
@app.route('/api/pdf-to-pptx', methods=['POST'])
def api_pdf_to_pptx():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded.'}), 400

    file = request.files['file']
    ext = os.path.splitext(file.filename)[1].lower()
    if ext != '.pdf':
        return jsonify({'error': 'Only PDF files are supported.'}), 400

    job_id = str(uuid.uuid4())
    job_upload = os.path.join(UPLOAD_DIR, job_id)
    job_output = os.path.join(OUTPUT_DIR, job_id)
    os.makedirs(job_upload, exist_ok=True)
    os.makedirs(job_output, exist_ok=True)

    try:
        from pptx import Presentation
        from pptx.util import Emu

        input_path = os.path.join(job_upload, file.filename)
        file.save(input_path)

        # Convert PDF pages to images
        images = convert_from_path(input_path, dpi=200)
        if not images:
            return jsonify({'error': 'Could not extract pages from PDF.'}), 500

        prs = Presentation()
        for i, img in enumerate(images):
            img_path = os.path.join(job_upload, f'page_{i+1}.png')
            img.save(img_path, 'PNG')

            # Set slide dimensions to match the image aspect ratio
            img_w, img_h = img.size
            slide_w = Emu(9144000)  # 10 inches in EMU
            slide_h = Emu(int(9144000 * img_h / img_w))
            prs.slide_width = slide_w
            prs.slide_height = slide_h

            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            slide.shapes.add_picture(img_path, Emu(0), Emu(0), slide_w, slide_h)

        base_name = os.path.splitext(file.filename)[0] + '.pptx'
        out_path = os.path.join(job_output, base_name)
        prs.save(out_path)

        if not os.path.exists(out_path):
            return jsonify({'error': 'Conversion produced no output.'}), 500

        response = send_file(
            out_path,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=base_name
        )

        @response.call_on_close
        def _cleanup():
            shutil.rmtree(job_upload, ignore_errors=True)
            shutil.rmtree(job_output, ignore_errors=True)

        return response

    except Exception as e:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': str(e)}), 500


# ─── PDF to Excel API ────────────────────────────────────────────────
@app.route('/api/pdf-to-excel', methods=['POST'])
def api_pdf_to_excel():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded.'}), 400

    file = request.files['file']
    ext = os.path.splitext(file.filename)[1].lower()
    if ext != '.pdf':
        return jsonify({'error': 'Only PDF files are supported.'}), 400

    job_id = str(uuid.uuid4())
    job_upload = os.path.join(UPLOAD_DIR, job_id)
    job_output = os.path.join(OUTPUT_DIR, job_id)
    os.makedirs(job_upload, exist_ok=True)
    os.makedirs(job_output, exist_ok=True)

    try:
        import pdfplumber
        from openpyxl import Workbook

        input_path = os.path.join(job_upload, file.filename)
        file.save(input_path)

        wb = Workbook()
        ws = wb.active
        ws.title = 'Sheet1'
        row_offset = 1

        with pdfplumber.open(input_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        for row in table:
                            clean_row = [(cell if cell else '') for cell in row]
                            ws.append(clean_row)
                        row_offset = ws.max_row + 1
                else:
                    # Fallback: extract text line-by-line
                    text = page.extract_text()
                    if text:
                        for line in text.split('\n'):
                            ws.append([line])

        base_name = os.path.splitext(file.filename)[0] + '.xlsx'
        out_path = os.path.join(job_output, base_name)
        wb.save(out_path)

        if not os.path.exists(out_path):
            return jsonify({'error': 'Conversion produced no output.'}), 500

        response = send_file(
            out_path,
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            as_attachment=True,
            download_name=base_name
        )

        @response.call_on_close
        def _cleanup():
            shutil.rmtree(job_upload, ignore_errors=True)
            shutil.rmtree(job_output, ignore_errors=True)

        return response

    except Exception as e:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': str(e)}), 500


# ─── PDF to CSV API ──────────────────────────────────────────────────
@app.route('/api/pdf-to-csv', methods=['POST'])
def api_pdf_to_csv():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded.'}), 400

    file = request.files['file']
    ext = os.path.splitext(file.filename)[1].lower()
    if ext != '.pdf':
        return jsonify({'error': 'Only PDF files are supported.'}), 400

    job_id = str(uuid.uuid4())
    job_upload = os.path.join(UPLOAD_DIR, job_id)
    job_output = os.path.join(OUTPUT_DIR, job_id)
    os.makedirs(job_upload, exist_ok=True)
    os.makedirs(job_output, exist_ok=True)

    try:
        import pdfplumber
        import csv

        input_path = os.path.join(job_upload, file.filename)
        file.save(input_path)

        base_name = os.path.splitext(file.filename)[0] + '.csv'
        out_path = os.path.join(job_output, base_name)

        with pdfplumber.open(input_path) as pdf:
            with open(out_path, 'w', newline='', encoding='utf-8') as f:
                writer = csv.writer(f)
                for page in pdf.pages:
                    tables = page.extract_tables()
                    if tables:
                        for table in tables:
                            for row in table:
                                clean_row = [(cell if cell else '') for cell in row]
                                writer.writerow(clean_row)
                    else:
                        # Fallback: extract text line-by-line
                        text = page.extract_text()
                        if text:
                            for line in text.split('\n'):
                                writer.writerow([line])

        if not os.path.exists(out_path) or os.path.getsize(out_path) == 0:
            return jsonify({'error': 'No data could be extracted from the PDF.'}), 500

        response = send_file(
            out_path,
            mimetype='text/csv',
            as_attachment=True,
            download_name=base_name
        )

        @response.call_on_close
        def _cleanup():
            shutil.rmtree(job_upload, ignore_errors=True)
            shutil.rmtree(job_output, ignore_errors=True)

        return response

    except Exception as e:
        shutil.rmtree(job_upload, ignore_errors=True)
        shutil.rmtree(job_output, ignore_errors=True)
        return jsonify({'error': str(e)}), 500


# ─── Page Count Helper ───────────────────────────────────────────────
@app.route('/api/pdf-page-count', methods=['POST'])
def api_pdf_page_count():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    file = request.files['file']
    if not file.filename.lower().endswith('.pdf'):
        return jsonify({'error': 'Only PDF files are allowed'}), 400

    job_id = str(uuid.uuid4())
    tmp_path = os.path.join(UPLOAD_DIR, f'{job_id}.pdf')
    file.save(tmp_path)
    try:
        reader = PdfReader(tmp_path)
        count = len(reader.pages)
        return jsonify({'pages': count})
    finally:
        os.remove(tmp_path)


# ─── Error Handlers ──────────────────────────────────────────────────
@app.errorhandler(413)
def too_large(e):
    return jsonify({'error': 'File too large. Maximum size is 50 MB.'}), 413


if __name__ == '__main__':
    debug = os.environ.get('FLASK_ENV') != 'production'
    app.run(host='0.0.0.0', debug=debug, port=5000)
